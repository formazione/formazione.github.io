from pptx import Presentation

# Open the PowerPoint presentation
prs = Presentation('presentation.pptx')

# Get the slide you want to remove
slide_to_remove = prs.slides[1]

# Remove the slide from the presentation
prs.slides.remove(slide_to_remove)

# Save the updated presentation
prs.save('updated_presentation.pptx')
